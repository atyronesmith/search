"""
Document extraction implementations
"""

import time
import json
import logging
from pathlib import Path
from typing import Dict, Any, Optional, List
import subprocess

import pypdfium2 as pdfium
from docx import Document
from pptx import Presentation

from models import DocumentElement, ExtractionResult, ElementType, BoundingBox

logger = logging.getLogger(__name__)


class BaseExtractor:
    """Base class for document extractors"""
    
    def __init__(self):
        self.extraction_method = "base"
    
    def extract(self, file_path: str, options: Optional[Dict[str, Any]] = None) -> ExtractionResult:
        """Extract content from document"""
        raise NotImplementedError


class PDFToTextExtractor(BaseExtractor):
    """PDF extraction using pdftotext (legacy fallback)"""
    
    def __init__(self):
        super().__init__()
        self.extraction_method = "pdftotext"
    
    def extract(self, file_path: str, options: Optional[Dict[str, Any]] = None) -> ExtractionResult:
        start_time = time.time()
        options = options or {}
        
        try:
            # Check if pdftotext is available
            result = subprocess.run(['which', 'pdftotext'], capture_output=True)
            if result.returncode != 0:
                return ExtractionResult(
                    success=False,
                    processing_time=time.time() - start_time,
                    error_message="pdftotext not available on system",
                    extraction_method=self.extraction_method,
                    metadata={'file_path': file_path}
                )
            
            # Run pdftotext
            cmd = ['pdftotext', '-layout', file_path, '-']
            if options.get('raw_text'):
                cmd = ['pdftotext', '-raw', file_path, '-']
            
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=30)
            processing_time = time.time() - start_time
            
            if result.returncode != 0:
                return ExtractionResult(
                    success=False,
                    processing_time=processing_time,
                    error_message=f"pdftotext failed: {result.stderr}",
                    extraction_method=self.extraction_method,
                    metadata={'file_path': file_path}
                )
            
            # Create text element
            elements = []
            if result.stdout.strip():
                elements.append(DocumentElement(
                    element_type=ElementType.PAGE_TEXT,
                    content=result.stdout.strip(),
                    page_number=0,
                    metadata={'extraction_method': self.extraction_method}
                ))
            
            return ExtractionResult(
                success=True,
                elements=elements,
                processing_time=processing_time,
                extraction_method=self.extraction_method,
                metadata={
                    'file_path': file_path,
                    'command_used': ' '.join(cmd)
                }
            )
            
        except subprocess.TimeoutExpired:
            return ExtractionResult(
                success=False,
                processing_time=time.time() - start_time,
                error_message="pdftotext timed out",
                extraction_method=self.extraction_method,
                metadata={'file_path': file_path}
            )
        except Exception as e:
            return ExtractionResult(
                success=False,
                processing_time=time.time() - start_time,
                error_message=f"pdftotext error: {str(e)}",
                extraction_method=self.extraction_method,
                metadata={'file_path': file_path}
            )


class PyPDFium2Extractor(BaseExtractor):
    """Enhanced PDF extraction using pypdfium2"""
    
    def __init__(self):
        super().__init__()
        self.extraction_method = "pypdfium2"
    
    def extract(self, file_path: str, options: Optional[Dict[str, Any]] = None) -> ExtractionResult:
        start_time = time.time()
        options = options or {}
        
        try:
            pdf = pdfium.PdfDocument(file_path)
            elements = []
            
            for page_num, page in enumerate(pdf):
                try:
                    textpage = page.get_textpage()
                    text = textpage.get_text_range()
                    
                    if text.strip():
                        # Basic structure detection
                        lines = text.strip().split('\n')
                        current_content = []
                        
                        for line in lines:
                            line = line.strip()
                            if not line:
                                continue
                                
                            # Simple heuristic for headings (all caps, short lines)
                            if len(line) < 100 and line.isupper() and len(line.split()) <= 8:
                                # Save previous content as paragraph
                                if current_content:
                                    elements.append(DocumentElement(
                                        element_type=ElementType.PARAGRAPH,
                                        content='\n'.join(current_content),
                                        page_number=page_num + 1,
                                        metadata={'extraction_method': self.extraction_method}
                                    ))
                                    current_content = []
                                
                                # Add as heading
                                elements.append(DocumentElement(
                                    element_type=ElementType.HEADING,
                                    content=line,
                                    page_number=page_num + 1,
                                    metadata={'extraction_method': self.extraction_method}
                                ))
                            else:
                                current_content.append(line)
                        
                        # Add remaining content
                        if current_content:
                            elements.append(DocumentElement(
                                element_type=ElementType.PARAGRAPH,
                                content='\n'.join(current_content),
                                page_number=page_num + 1,
                                metadata={'extraction_method': self.extraction_method}
                            ))
                    
                    textpage.close()
                except Exception as e:
                    logger.warning(f"Error processing page {page_num + 1}: {e}")
                finally:
                    page.close()
            
            pdf.close()
            processing_time = time.time() - start_time
            
            return ExtractionResult(
                success=True,
                elements=elements,
                processing_time=processing_time,
                extraction_method=self.extraction_method,
                metadata={
                    'file_path': file_path,
                    'total_pages': page_num + 1 if 'page_num' in locals() else 0
                }
            )
            
        except Exception as e:
            return ExtractionResult(
                success=False,
                processing_time=time.time() - start_time,
                error_message=f"pypdfium2 error: {str(e)}",
                extraction_method=self.extraction_method,
                metadata={'file_path': file_path}
            )


class DocxExtractor(BaseExtractor):
    """DOCX document extraction using python-docx"""
    
    def __init__(self):
        super().__init__()
        self.extraction_method = "python-docx"
    
    def extract(self, file_path: str, options: Optional[Dict[str, Any]] = None) -> ExtractionResult:
        start_time = time.time()
        options = options or {}
        
        try:
            doc = Document(file_path)
            elements = []
            
            # Extract paragraphs with style information
            for para in doc.paragraphs:
                if para.text.strip():
                    # Determine element type based on style
                    element_type = ElementType.PARAGRAPH
                    if para.style.name.startswith('Heading'):
                        element_type = ElementType.HEADING
                    elif para.style.name in ['Title', 'Subtitle']:
                        element_type = ElementType.TITLE
                    
                    elements.append(DocumentElement(
                        element_type=element_type,
                        content=para.text.strip(),
                        metadata={
                            'style': para.style.name,
                            'extraction_method': self.extraction_method,
                            'paragraph_alignment': str(para.alignment) if para.alignment else None
                        }
                    ))
            
            # Extract tables
            for table_idx, table in enumerate(doc.tables):
                table_content = []
                for row in table.rows:
                    row_content = [cell.text.strip() for cell in row.cells]
                    table_content.append(row_content)
                
                # Convert table to JSON string for storage
                table_json = json.dumps(table_content, ensure_ascii=False)
                
                elements.append(DocumentElement(
                    element_type=ElementType.TABLE,
                    content=table_json,
                    metadata={
                        'table_index': table_idx,
                        'rows': len(table.rows),
                        'columns': len(table.columns),
                        'extraction_method': self.extraction_method
                    }
                ))
            
            processing_time = time.time() - start_time
            
            return ExtractionResult(
                success=True,
                elements=elements,
                processing_time=processing_time,
                extraction_method=self.extraction_method,
                metadata={
                    'file_path': file_path,
                    'paragraph_count': len([e for e in elements if e.element_type in [ElementType.PARAGRAPH, ElementType.HEADING, ElementType.TITLE]]),
                    'table_count': len([e for e in elements if e.element_type == ElementType.TABLE])
                }
            )
            
        except Exception as e:
            return ExtractionResult(
                success=False,
                processing_time=time.time() - start_time,
                error_message=f"python-docx error: {str(e)}",
                extraction_method=self.extraction_method,
                metadata={'file_path': file_path}
            )


class PptxExtractor(BaseExtractor):
    """PPTX presentation extraction using python-pptx"""
    
    def __init__(self):
        super().__init__()
        self.extraction_method = "python-pptx"
    
    def extract(self, file_path: str, options: Optional[Dict[str, Any]] = None) -> ExtractionResult:
        start_time = time.time()
        options = options or {}
        
        try:
            prs = Presentation(file_path)
            elements = []
            
            for slide_num, slide in enumerate(prs.slides):
                # Process slide title
                if slide.shapes.title and slide.shapes.title.text.strip():
                    elements.append(DocumentElement(
                        element_type=ElementType.SLIDE_TITLE,
                        content=slide.shapes.title.text.strip(),
                        page_number=slide_num + 1,
                        metadata={
                            'slide_number': slide_num + 1,
                            'extraction_method': self.extraction_method
                        }
                    ))
                
                # Process other shapes
                for shape_idx, shape in enumerate(slide.shapes):
                    if hasattr(shape, 'text') and shape.text.strip():
                        # Skip title shape as it's already processed
                        if shape == slide.shapes.title:
                            continue
                        
                        # Determine element type
                        element_type = ElementType.SLIDE_CONTENT
                        
                        # Check if it's a table
                        if hasattr(shape, 'table'):
                            element_type = ElementType.TABLE
                            # Extract table data
                            table_data = []
                            for row in shape.table.rows:
                                row_data = [cell.text.strip() for cell in row.cells]
                                table_data.append(row_data)
                            content = json.dumps(table_data, ensure_ascii=False)
                        else:
                            content = shape.text.strip()
                        
                        elements.append(DocumentElement(
                            element_type=element_type,
                            content=content,
                            page_number=slide_num + 1,
                            metadata={
                                'slide_number': slide_num + 1,
                                'shape_index': shape_idx,
                                'shape_type': str(shape.shape_type) if hasattr(shape, 'shape_type') else None,
                                'extraction_method': self.extraction_method
                            }
                        ))
            
            processing_time = time.time() - start_time
            
            return ExtractionResult(
                success=True,
                elements=elements,
                processing_time=processing_time,
                extraction_method=self.extraction_method,
                metadata={
                    'file_path': file_path,
                    'slide_count': len(prs.slides),
                    'title_count': len([e for e in elements if e.element_type == ElementType.SLIDE_TITLE]),
                    'content_count': len([e for e in elements if e.element_type == ElementType.SLIDE_CONTENT])
                }
            )
            
        except Exception as e:
            return ExtractionResult(
                success=False,
                processing_time=time.time() - start_time,
                error_message=f"python-pptx error: {str(e)}",
                extraction_method=self.extraction_method,
                metadata={'file_path': file_path}
            )


class DocumentExtractorManager:
    """Manages different document extractors"""
    
    def __init__(self):
        self.extractors = {
            'pdftotext': PDFToTextExtractor(),
            'pypdfium2': PyPDFium2Extractor(),
            'python-docx': DocxExtractor(),
            'python-pptx': PptxExtractor()
        }
    
    def get_extractor_for_file(self, file_path: str, preferred_method: Optional[str] = None) -> Optional[BaseExtractor]:
        """Get appropriate extractor for file type"""
        path = Path(file_path)
        extension = path.suffix.lower()
        
        # If specific method requested, try to use it
        if preferred_method and preferred_method in self.extractors:
            return self.extractors[preferred_method]
        
        # Auto-detect based on file extension
        if extension == '.pdf':
            return self.extractors['pypdfium2']  # Prefer pypdfium2 over pdftotext
        elif extension == '.docx':
            return self.extractors['python-docx']
        elif extension == '.pptx':
            return self.extractors['python-pptx']
        else:
            return None
    
    def extract_document(self, file_path: str, extraction_method: str = "auto", 
                        options: Optional[Dict[str, Any]] = None) -> ExtractionResult:
        """Extract document using appropriate method"""
        if not Path(file_path).exists():
            return ExtractionResult(
                success=False,
                processing_time=0.0,
                error_message=f"File not found: {file_path}",
                extraction_method="none",
                metadata={'file_path': file_path}
            )
        
        extractor = self.get_extractor_for_file(file_path, extraction_method)
        if not extractor:
            return ExtractionResult(
                success=False,
                processing_time=0.0,
                error_message=f"No suitable extractor found for file: {file_path}",
                extraction_method="none",
                metadata={'file_path': file_path}
            )
        
        return extractor.extract(file_path, options)