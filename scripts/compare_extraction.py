#!/usr/bin/env python3
"""
Document Extraction Comparison Script

Compares different document processing methods:
1. Current pdftotext approach (baseline)
2. pypdfium2 for enhanced PDF processing
3. python-docx for DOCX files
4. python-pptx for PPTX files
5. Future: docling integration (when dependencies resolve)
"""

import subprocess
import time
import json
import argparse
from pathlib import Path
from typing import Dict, Any, Optional, List
import pypdfium2 as pdfium
from docx import Document
from pptx import Presentation
import logging

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class DocumentElement:
    """Represents a structured document element"""
    def __init__(self, element_type: str, content: str, page_number: int = 0, 
                 metadata: Optional[Dict] = None):
        self.element_type = element_type
        self.content = content
        self.page_number = page_number
        self.metadata = metadata or {}


class ExtractionResult:
    """Contains extraction results and metadata"""
    def __init__(self, success: bool, elements: List[DocumentElement] = None, 
                 processing_time: float = 0.0, error_message: Optional[str] = None):
        self.success = success
        self.elements = elements or []
        self.processing_time = processing_time
        self.error_message = error_message
        
    def to_dict(self):
        return {
            'success': self.success,
            'elements': [{
                'type': elem.element_type,
                'content': elem.content[:200] + '...' if len(elem.content) > 200 else elem.content,
                'page_number': elem.page_number,
                'metadata': elem.metadata
            } for elem in self.elements],
            'element_count': len(self.elements),
            'processing_time': self.processing_time,
            'error_message': self.error_message
        }


class PDFToTextExtractor:
    """Current pdftotext extraction method"""
    
    def extract(self, pdf_path: str) -> ExtractionResult:
        start_time = time.time()
        try:
            result = subprocess.run(
                ['pdftotext', '-layout', pdf_path, '-'], 
                capture_output=True, text=True, timeout=30
            )
            processing_time = time.time() - start_time
            
            if result.returncode != 0:
                return ExtractionResult(
                    success=False,
                    processing_time=processing_time,
                    error_message=f"pdftotext failed: {result.stderr}"
                )
            
            # Create simple text element
            elements = [DocumentElement(
                element_type='text',
                content=result.stdout.strip(),
                page_number=0
            )]
            
            return ExtractionResult(
                success=True,
                elements=elements,
                processing_time=processing_time
            )
            
        except subprocess.TimeoutExpired:
            return ExtractionResult(
                success=False,
                processing_time=time.time() - start_time,
                error_message="pdftotext timed out"
            )
        except Exception as e:
            return ExtractionResult(
                success=False,
                processing_time=time.time() - start_time,
                error_message=f"pdftotext error: {str(e)}"
            )


class PyPDFium2Extractor:
    """Enhanced PDF extraction using pypdfium2"""
    
    def extract(self, pdf_path: str) -> ExtractionResult:
        start_time = time.time()
        try:
            pdf = pdfium.PdfDocument(pdf_path)
            elements = []
            
            for page_num, page in enumerate(pdf):
                textpage = page.get_textpage()
                text = textpage.get_text_range()
                
                if text.strip():
                    elements.append(DocumentElement(
                        element_type='page_text',
                        content=text.strip(),
                        page_number=page_num + 1,
                        metadata={'extraction_method': 'pypdfium2'}
                    ))
                
                textpage.close()
                page.close()
            
            pdf.close()
            processing_time = time.time() - start_time
            
            return ExtractionResult(
                success=True,
                elements=elements,
                processing_time=processing_time
            )
            
        except Exception as e:
            return ExtractionResult(
                success=False,
                processing_time=time.time() - start_time,
                error_message=f"pypdfium2 error: {str(e)}"
            )


class DocxExtractor:
    """DOCX document extraction"""
    
    def extract(self, docx_path: str) -> ExtractionResult:
        start_time = time.time()
        try:
            doc = Document(docx_path)
            elements = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    # Determine element type based on style
                    element_type = 'paragraph'
                    if para.style.name.startswith('Heading'):
                        element_type = 'heading'
                    elif para.style.name == 'Title':
                        element_type = 'title'
                    
                    elements.append(DocumentElement(
                        element_type=element_type,
                        content=para.text.strip(),
                        metadata={
                            'style': para.style.name,
                            'extraction_method': 'python-docx'
                        }
                    ))
            
            # Extract tables
            for table in doc.tables:
                table_content = []
                for row in table.rows:
                    row_content = [cell.text.strip() for cell in row.cells]
                    table_content.append(row_content)
                
                elements.append(DocumentElement(
                    element_type='table',
                    content=json.dumps(table_content),
                    metadata={
                        'rows': len(table.rows),
                        'columns': len(table.columns),
                        'extraction_method': 'python-docx'
                    }
                ))
            
            processing_time = time.time() - start_time
            
            return ExtractionResult(
                success=True,
                elements=elements,
                processing_time=processing_time
            )
            
        except Exception as e:
            return ExtractionResult(
                success=False,
                processing_time=time.time() - start_time,
                error_message=f"python-docx error: {str(e)}"
            )


class PptxExtractor:
    """PPTX presentation extraction"""
    
    def extract(self, pptx_path: str) -> ExtractionResult:
        start_time = time.time()
        try:
            prs = Presentation(pptx_path)
            elements = []
            
            for slide_num, slide in enumerate(prs.slides):
                slide_content = []
                
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        # Determine if it's a title or content
                        element_type = 'slide_content'
                        if shape.shape_type == 1:  # Title placeholder
                            element_type = 'slide_title'
                        
                        elements.append(DocumentElement(
                            element_type=element_type,
                            content=shape.text.strip(),
                            page_number=slide_num + 1,
                            metadata={
                                'slide_number': slide_num + 1,
                                'shape_type': str(shape.shape_type),
                                'extraction_method': 'python-pptx'
                            }
                        ))
            
            processing_time = time.time() - start_time
            
            return ExtractionResult(
                success=True,
                elements=elements,
                processing_time=processing_time
            )
            
        except Exception as e:
            return ExtractionResult(
                success=False,
                processing_time=time.time() - start_time,
                error_message=f"python-pptx error: {str(e)}"
            )


def compare_extraction(file_path: str) -> Dict[str, Any]:
    """Compare extraction methods for a given document"""
    path = Path(file_path)
    
    if not path.exists():
        return {'error': f'File not found: {file_path}'}
    
    results = {
        'file_path': file_path,
        'file_size': path.stat().st_size,
        'file_extension': path.suffix.lower(),
        'extractions': {}
    }
    
    # PDF extraction comparison
    if path.suffix.lower() == '.pdf':
        logger.info(f"Testing PDF extraction methods for {path.name}")
        
        # Method 1: pdftotext (current)
        pdftotext_extractor = PDFToTextExtractor()
        results['extractions']['pdftotext'] = pdftotext_extractor.extract(file_path).to_dict()
        
        # Method 2: pypdfium2 (enhanced)
        pypdfium2_extractor = PyPDFium2Extractor()
        results['extractions']['pypdfium2'] = pypdfium2_extractor.extract(file_path).to_dict()
    
    # DOCX extraction
    elif path.suffix.lower() == '.docx':
        logger.info(f"Testing DOCX extraction for {path.name}")
        docx_extractor = DocxExtractor()
        results['extractions']['python-docx'] = docx_extractor.extract(file_path).to_dict()
    
    # PPTX extraction
    elif path.suffix.lower() == '.pptx':
        logger.info(f"Testing PPTX extraction for {path.name}")
        pptx_extractor = PptxExtractor()
        results['extractions']['python-pptx'] = pptx_extractor.extract(file_path).to_dict()
    
    else:
        results['extractions']['unsupported'] = {
            'success': False,
            'error_message': f'Unsupported file type: {path.suffix}'
        }
    
    return results


def main():
    parser = argparse.ArgumentParser(description='Compare document extraction methods')
    parser.add_argument('files', nargs='+', help='Document files to process')
    parser.add_argument('--output', '-o', help='Output JSON file')
    parser.add_argument('--verbose', '-v', action='store_true', help='Verbose output')
    
    args = parser.parse_args()
    
    if args.verbose:
        logging.getLogger().setLevel(logging.DEBUG)
    
    all_results = []
    
    for file_path in args.files:
        logger.info(f"Processing: {file_path}")
        result = compare_extraction(file_path)
        all_results.append(result)
        
        if args.verbose:
            print(f"\nResults for {file_path}:")
            print(json.dumps(result, indent=2))
    
    # Summary statistics
    summary = {
        'total_files': len(all_results),
        'successful_extractions': 0,
        'failed_extractions': 0,
        'average_processing_time': 0.0,
        'files_by_type': {},
        'results': all_results
    }
    
    total_time = 0.0
    extraction_count = 0
    
    for result in all_results:
        file_ext = result.get('file_extension', 'unknown')
        if file_ext not in summary['files_by_type']:
            summary['files_by_type'][file_ext] = 0
        summary['files_by_type'][file_ext] += 1
        
        for method, extraction in result.get('extractions', {}).items():
            extraction_count += 1
            if extraction.get('success', False):
                summary['successful_extractions'] += 1
                total_time += extraction.get('processing_time', 0.0)
            else:
                summary['failed_extractions'] += 1
    
    if extraction_count > 0:
        summary['average_processing_time'] = total_time / extraction_count
    
    if args.output:
        with open(args.output, 'w') as f:
            json.dump(summary, f, indent=2)
        logger.info(f"Results saved to {args.output}")
    else:
        print(json.dumps(summary, indent=2))


if __name__ == '__main__':
    main()